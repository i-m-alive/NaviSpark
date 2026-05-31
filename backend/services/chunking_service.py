"""
Chunking service for large PDF and PPTX documents.

When a document exceeds CHUNK_THRESHOLD pages/slides, this service:
  1. Splits it into overlapping chunks (CHUNK_SIZE pages, CHUNK_OVERLAP overlap)
  2. Sends each chunk to Bedrock for structured summarisation (in parallel)
  3. Merges all chunk summaries into a single unified context JSON string
  4. Returns that string to callers — agents receive it instead of the raw file

For PDF: chunks are standalone PDF bytes sent to Bedrock as document blocks.
For PPTX/PPT: chunks are slide-text strings sent as plain text (consistent with
  how bedrock_client already handles PPTX files in invoke_agent_with_pdf).

Thresholds (overridable via env vars):
  CHUNK_THRESHOLD = 30  — documents with more pages/slides are chunked
  CHUNK_SIZE      = 15  — pages/slides per chunk
  CHUNK_OVERLAP   = 3   — overlap pages between adjacent chunks
"""

import io
import json
import logging
import os
import time

from pypdf import PdfReader, PdfWriter
from pptx import Presentation

logger = logging.getLogger(__name__)

CHUNK_THRESHOLD: int = int(os.getenv("CHUNK_THRESHOLD", "30"))
CHUNK_SIZE: int = int(os.getenv("CHUNK_SIZE", "15"))
CHUNK_OVERLAP: int = int(os.getenv("CHUNK_OVERLAP", "3"))


# ── PDF Splitting ──────────────────────────────────────────────────────────────

def split_pdf_into_chunks(
    pdf_bytes: bytes,
    chunk_size: int = CHUNK_SIZE,
    overlap: int = CHUNK_OVERLAP,
) -> list:
    """
    Splits a PDF into overlapping page-range chunks using pypdf.

    Returns a list of (chunk_pdf_bytes, start_page_1based, end_page_1based) tuples,
    ordered from first to last chunk. Each chunk is a fully valid standalone PDF.

    Example — 45-page doc, chunk_size=15, overlap=3:
      Chunk 0: pages  1-15
      Chunk 1: pages 13-27  (overlap: pages 13-15 appear in both chunk 0 and chunk 1)
      Chunk 2: pages 25-39
      Chunk 3: pages 37-45  (final chunk may be shorter than chunk_size)
    """
    reader = PdfReader(io.BytesIO(pdf_bytes))
    total_pages = len(reader.pages)
    chunks = []
    start = 0

    while start < total_pages:
        end = min(start + chunk_size, total_pages)

        writer = PdfWriter()
        for page_num in range(start, end):
            writer.add_page(reader.pages[page_num])

        buf = io.BytesIO()
        writer.write(buf)
        chunks.append((buf.getvalue(), start + 1, end))

        if end >= total_pages:
            break
        start = end - overlap

    return chunks


# ── PPTX Splitting ─────────────────────────────────────────────────────────────

def split_pptx_into_text_chunks(
    pptx_bytes: bytes,
    chunk_size: int = CHUNK_SIZE,
    overlap: int = CHUNK_OVERLAP,
) -> list:
    """
    Splits a PPTX/PPT into overlapping slide-text chunks.

    Returns a list of (slide_text_string, start_slide_1based, end_slide_1based) tuples.
    PPTX files are handled as extracted text (consistent with bedrock_client's
    existing _extract_pptx_text behaviour — no binary PPTX splitting needed).
    """
    prs = Presentation(io.BytesIO(pptx_bytes))
    slides = list(prs.slides)
    total_slides = len(slides)
    chunks = []
    start = 0

    while start < total_slides:
        end = min(start + chunk_size, total_slides)

        slide_texts = []
        for i in range(start, end):
            texts = []
            for shape in slides[i].shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
            if texts:
                slide_texts.append(f"--- Slide {i + 1} ---\n" + "\n".join(texts))

        chunk_text = (
            "\n\n".join(slide_texts)
            if slide_texts
            else "(No text content in this slide range)"
        )
        chunks.append((chunk_text, start + 1, end))

        if end >= total_slides:
            break
        start = end - overlap

    return chunks


# ── Chunk Summarisation ────────────────────────────────────────────────────────

# Gap between consecutive chunk summarisation calls (seconds).
# Sending chunks serially with a small gap prevents bursting the Bedrock TPM
# quota on dev/hackathon AWS accounts (ThrottlingException).
_CHUNK_SERIAL_DELAY_SECS: float = float(os.getenv("CHUNK_SERIAL_DELAY", "3"))


def _make_fallback_summary(chunk_index: int, start: int, end: int, error: str) -> dict:
    """Minimal summary used when a chunk's Bedrock call fails."""
    return {
        "chunk_index": chunk_index,
        "page_range": {"start": start, "end": end},
        "summary": f"[Chunk {chunk_index} summarisation failed: {error}]",
        "sections_found": [],
        "key_claims": [],
        "figures_and_numbers": [],
        "scope_items": [],
        "risks": [],
        "commercial_terms": [],
    }


def _summarize_chunks_serial(
    chunks: list,
    file_type: str,
    sid: str = "????????",
    emit=None,
) -> list:
    """
    Runs chunk summarisations one at a time with a short sleep between calls.
    sid is a short session ID used for log line prefixes.
    emit is an optional callable emit(activity, status) for the activity feed.
    """
    from bedrock_client import summarize_chunk_pdf, summarize_chunk_pptx_text

    _e = emit if emit else (lambda a, s="running": None)
    total = len(chunks)
    results = []

    for idx, chunk in enumerate(chunks):
        data, start, end = chunk

        if idx > 0:
            logger.info("[CHUNKING] [%s] Waiting %ss before next chunk call (rate-limit gap)...",
                        sid, _CHUNK_SERIAL_DELAY_SECS)
            time.sleep(_CHUNK_SERIAL_DELAY_SECS)

        logger.info("[CHUNKING] [%s] Chunk %d/%d — pages %d-%d — sending to Bedrock...",
                    sid, idx + 1, total, start, end)
        _e(f"Summarizing chunk {idx + 1}/{total} (pages {start}–{end})")
        t0 = time.monotonic()

        try:
            if file_type == "pdf":
                summary = summarize_chunk_pdf(data, idx, start, end)
            else:
                summary = summarize_chunk_pptx_text(data, idx, start, end)
            elapsed = time.monotonic() - t0
            results.append(summary)
            logger.info("[CHUNKING] [%s] Chunk %d/%d — pages %d-%d — DONE (%.1fs)",
                        sid, idx + 1, total, start, end, elapsed)
            _e(f"Chunk {idx + 1}/{total} summarized (pages {start}–{end})", "completed")
        except Exception as exc:
            elapsed = time.monotonic() - t0
            logger.error("[CHUNKING] [%s] Chunk %d/%d — pages %d-%d — FAILED after %.1fs: %s",
                         sid, idx + 1, total, start, end, elapsed, exc)
            _e(f"Chunk {idx + 1}/{total} failed: {exc}", "error")
            results.append(_make_fallback_summary(idx, start, end, str(exc)))

    return results


# ── Merge ──────────────────────────────────────────────────────────────────────

def _merge_summaries(
    summaries: list,
    total_pages: int,
    client_industry: list,
    proposal_type: str,
    client_priorities: list,
) -> str:
    """
    Merges all chunk summary dicts into a single unified context JSON string.
    This string is what Agents 1, 2, and 3 receive instead of the raw document
    when the original file exceeded CHUNK_THRESHOLD pages/slides.
    """
    context = {
        "document_type": "chunked_proposal_context",
        "full_page_count": total_pages,
        "total_chunks": len(summaries),
        "client_context": {
            "industry": client_industry,
            "proposal_type": proposal_type,
            "priorities": client_priorities,
        },
        "instruction": (
            "This context was produced by splitting a large document into overlapping chunks "
            "and summarising each chunk via Claude. All chunks together represent the complete "
            "proposal. Review ALL entries in document_summary — do not skip any chunk."
        ),
        "document_summary": summaries,
    }
    return json.dumps(context, indent=2)


# ── Main Entry Point ───────────────────────────────────────────────────────────

def prepare_document_context(
    file_bytes: bytes,
    file_type: str,
    page_count: int,
    client_industry: list,
    proposal_type: str,
    client_priorities: list,
    sid: str = "????????",
    emit=None,
) -> str | None:
    """
    Decides whether chunking is needed for this document.

    If page_count > CHUNK_THRESHOLD:
      1. Splits the document into overlapping chunks
      2. Summarises each chunk serially via Bedrock
      3. Merges summaries into a unified context JSON string
      4. Returns the context string — callers pass this to agent run() as
         pre_processed_context instead of raw file bytes

    If page_count <= CHUNK_THRESHOLD:
      Returns None — the caller should use the raw file (current behaviour).

    sid is a short session ID used for log line prefixes.
    """
    _e = emit if emit else (lambda a, s="running": None)

    if page_count <= CHUNK_THRESHOLD:
        logger.info("[CHUNKING] [%s] %d pages <= threshold %d — skipping chunking, raw file goes to agents.",
                    sid, page_count, CHUNK_THRESHOLD)
        return None

    logger.info("─" * 60)
    logger.info("[CHUNKING] [%s] %d pages > threshold %d — ENTERING chunking pipeline",
                sid, page_count, CHUNK_THRESHOLD)
    logger.info("[CHUNKING] [%s] Settings: chunk_size=%d pages, overlap=%d pages",
                sid, CHUNK_SIZE, CHUNK_OVERLAP)

    _e(f"Document has {page_count} pages — activating chunking pipeline")

    t_split = time.monotonic()
    if file_type == "pdf":
        chunks = split_pdf_into_chunks(file_bytes)
    else:
        chunks = split_pptx_into_text_chunks(file_bytes)

    logger.info("[CHUNKING] [%s] Split complete (%.2fs) — %d chunks created:",
                sid, time.monotonic() - t_split, len(chunks))
    for i, (_, s, e) in enumerate(chunks):
        logger.info("[CHUNKING] [%s]   Chunk %d: pages %d-%d (%d pages)", sid, i + 1, s, e, e - s + 1)

    _e(f"Split into {len(chunks)} chunks ({CHUNK_SIZE} pages each)", "completed")

    logger.info("[CHUNKING] [%s] Starting serial summarisation — %d Bedrock calls, ~%ds apart...",
                sid, len(chunks), int(_CHUNK_SERIAL_DELAY_SECS))
    _e(f"Summarizing {len(chunks)} chunks via AI (rate-limited)")

    t_summarise = time.monotonic()
    summaries = _summarize_chunks_serial(chunks, file_type, sid, emit=_e)
    logger.info("[CHUNKING] [%s] All %d chunks summarised in %.1fs — merging context...",
                sid, len(summaries), time.monotonic() - t_summarise)

    _e("Merging chunk summaries into unified context")
    merged = _merge_summaries(summaries, page_count, client_industry, proposal_type, client_priorities)
    logger.info("[CHUNKING] [%s] Merged context ready — %.1f KB  (agents will receive this as text)",
                sid, len(merged) / 1024)
    logger.info("─" * 60)
    _e(f"Document pre-processing complete ({len(merged) / 1024:.0f} KB context)", "completed")

    return merged
