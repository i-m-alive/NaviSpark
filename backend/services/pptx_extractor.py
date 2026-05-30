"""
Extracts structured text content from a PPTX/PPT file.
Returns a slide map that Agent 5 uses to locate where each finding lives.
"""

import io
from typing import Optional
from pptx import Presentation


def _classify_shape(shape) -> str:
    """Return a simple content-role label for a shape."""
    if not shape.is_placeholder:
        return "free_text"
    ph = shape.placeholder_format
    if ph is None:
        return "free_text"
    if ph.idx == 0:
        return "title"
    if ph.idx == 1:
        return "body"
    return "other_placeholder"


def _max_font_pt(shape) -> Optional[int]:
    """Return the largest explicit font size (in pt) found in any run, or None."""
    sizes = []
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if run.font.size:
                sizes.append(int(run.font.size.pt))
    return max(sizes) if sizes else None


def extract_slide_map(pptx_bytes: bytes) -> list:
    """
    Parses a PPTX and returns a structured list of slides with shape content.

    Each slide entry:
      {
        "slide_index": int,           # 0-based
        "slide_title": str,           # best-guess title
        "shapes": [
          {
            "shape_name": str,        # exact name as stored in PPTX (use this in modifications)
            "shape_index": int,       # 1-based shape ID
            "shape_role": str,        # "title" | "body" | "other_placeholder" | "free_text"
            "is_placeholder": bool,
            "placeholder_idx": int | None,  # 0=title, 1=body, etc.
            "max_font_pt": int | None,      # largest explicit font size; None = inherited
            "text": str,              # full concatenated text
            "paragraphs": [str],      # each non-empty paragraph separately
          }
        ]
      }
    """
    prs = Presentation(io.BytesIO(pptx_bytes))
    slides = []

    for slide_index, slide in enumerate(prs.slides):
        slide_title = ""
        shapes = []

        for shape in slide.shapes:
            # Skip non-text shapes (images, charts, etc.)
            if not shape.has_text_frame:
                continue

            tf = shape.text_frame
            full_text = tf.text.strip()
            if not full_text:
                continue

            # Extract each non-empty paragraph separately
            paragraphs = [p.text.strip() for p in tf.paragraphs if p.text.strip()]

            # Identify placeholder type
            placeholder_idx = None
            is_placeholder = shape.is_placeholder
            if is_placeholder and shape.placeholder_format is not None:
                placeholder_idx = shape.placeholder_format.idx

            # Title is placeholder index 0
            if placeholder_idx == 0:
                slide_title = full_text

            # Skip large-font non-placeholder shapes: these are design/decorative
            # elements (display callouts, watermarks, oversized quotes).  Exposing
            # them to Agent 5 causes it to write replacement text into them, which
            # visually overlays the real body content.  28 pt is a safe threshold:
            # body text is ≤18 pt; genuine content headings use placeholders.
            max_fp = _max_font_pt(shape)
            if max_fp is not None and max_fp > 28 and not is_placeholder:
                continue

            shapes.append({
                "shape_name": shape.name,
                "shape_index": shape.shape_id,
                "shape_role": _classify_shape(shape),
                "is_placeholder": is_placeholder,
                "placeholder_idx": placeholder_idx,
                "max_font_pt": _max_font_pt(shape),
                "text": full_text,
                "paragraphs": paragraphs,
            })

        slides.append({
            "slide_index": slide_index,
            "slide_title": slide_title or f"Slide {slide_index + 1}",
            "shapes": shapes,
        })

    return slides


def slide_map_to_prompt_text(slide_map: list) -> str:
    """
    Converts the slide map to a compact, readable text block for the Agent 5 prompt.

    Format per shape:
        [EXACT_SHAPE_NAME] (role, font) : text content

    The text inside [...] is the EXACT shape_name to use in modification JSON.
    The (...) parenthetical is metadata only — never include it in shape_name.
    Roles: title | body | other_placeholder | free_text
    Font: explicit max pt if set in the file; "inherited" means it follows the slide master.
    Free-text shapes with large explicit fonts are design/callout elements — avoid modifying them.
    """
    lines = []
    for slide in slide_map:
        idx = slide["slide_index"]
        title = slide["slide_title"]
        lines.append(f"\n=== Slide {idx} — {title} ===")
        for shape in slide["shapes"]:
            name = shape["shape_name"]
            role = shape.get("shape_role", "free_text")
            font_pt = shape.get("max_font_pt")
            font_label = f"{font_pt}pt" if font_pt else "inherited"
            text = shape["text"]
            lines.append(
                f'  [{name}] ({role}, {font_label}): {text[:500]}{"..." if len(text) > 500 else ""}'
            )
    return "\n".join(lines)
