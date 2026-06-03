"""
Utility helpers for the Custom Checklist Review Pipeline.

  extract_proposal_text  — PDF/PPTX bytes → plain text for NC1
  download_checklist_to_tempfile — downloads checklist from Supabase Storage to a
                                   named temp file so NC2 can read it by path
  checklist_ext_from_filename    — validates and returns the file extension
"""

import io
import logging
import os
import tempfile

logger = logging.getLogger(__name__)

SUPPORTED_CHECKLIST_EXTS = {".xlsx", ".xlsm", ".csv", ".docx", ".pdf"}


# ── Proposal text extraction ──────────────────────────────────────────────────

def extract_proposal_text(file_bytes: bytes, file_type: str) -> str:
    """Extract plain text from a proposal PDF or PPTX for NC1 consumption."""
    if file_type == "pdf":
        return _extract_pdf_text(file_bytes)
    if file_type in ("pptx", "ppt"):
        return _extract_pptx_text(file_bytes)
    raise ValueError(f"Unsupported proposal file type for text extraction: '{file_type}'")


def _extract_pdf_text(pdf_bytes: bytes) -> str:
    try:
        from pypdf import PdfReader
    except ImportError:
        try:
            from PyPDF2 import PdfReader  # type: ignore[no-redef]
        except ImportError:
            raise ImportError("pypdf (or PyPDF2) is required for PDF text extraction")

    reader = PdfReader(io.BytesIO(pdf_bytes))
    parts: list[str] = []
    for page_num, page in enumerate(reader.pages):
        try:
            text = page.extract_text() or ""
            stripped = text.strip()
            if stripped:
                parts.append(f"[Page {page_num + 1}]\n{stripped}")
        except Exception as exc:
            logger.warning("PDF page %d extraction failed (skipped): %s", page_num + 1, exc)

    full_text = "\n\n".join(parts)
    logger.info("PDF text extracted: %d pages, %d chars", len(reader.pages), len(full_text))
    return full_text


def _extract_pptx_text(pptx_bytes: bytes) -> str:
    try:
        from pptx import Presentation  # type: ignore[import]
    except ImportError:
        raise ImportError("python-pptx is required for PPTX text extraction")

    prs = Presentation(io.BytesIO(pptx_bytes))
    parts: list[str] = []
    for i, slide in enumerate(prs.slides):
        slide_lines: list[str] = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    slide_lines.append(text)
        if slide_lines:
            parts.append(f"[Slide {i + 1}]\n" + "\n".join(slide_lines))

    full_text = "\n\n".join(parts)
    logger.info("PPTX text extracted: %d slides, %d chars", len(prs.slides), len(full_text))
    return full_text


# ── Checklist file handling ───────────────────────────────────────────────────

def checklist_ext_from_filename(filename: str) -> str:
    """Return the validated lowercase extension (e.g. '.xlsx') or raise ValueError."""
    _, ext = os.path.splitext(filename.lower())
    if ext not in SUPPORTED_CHECKLIST_EXTS:
        raise ValueError(
            f"Unsupported checklist format '{ext}'. "
            "Supported: Excel (.xlsx/.xlsm), CSV (.csv), Word (.docx), PDF (.pdf)"
        )
    return ext


def download_checklist_to_tempfile(checklist_storage_path: str) -> str:
    """
    Download a checklist file from Supabase Storage and write it to a temporary
    file on disk. Returns the absolute path of the temp file.

    The caller is responsible for deleting the temp file after use:
        os.unlink(tmp_path)
    """
    from storage import download_file_from_storage  # local import avoids circular deps

    file_bytes = download_file_from_storage(checklist_storage_path)
    if not file_bytes:
        raise ValueError(f"Downloaded checklist is empty: {checklist_storage_path}")

    ext = os.path.splitext(checklist_storage_path)[1].lower() or ".tmp"
    if ext not in SUPPORTED_CHECKLIST_EXTS:
        raise ValueError(
            f"Checklist file extension '{ext}' is not supported. "
            f"Supported: {', '.join(sorted(SUPPORTED_CHECKLIST_EXTS))}"
        )

    fd, tmp_path = tempfile.mkstemp(suffix=ext, prefix="nc2_checklist_")
    try:
        with os.fdopen(fd, "wb") as fh:
            fh.write(file_bytes)
    except Exception:
        try:
            os.close(fd)
        except Exception:
            pass
        raise

    logger.info(
        "Checklist written to temp file: %s  (%.1f KB)",
        tmp_path, len(file_bytes) / 1024,
    )
    return tmp_path
