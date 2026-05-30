"""
Applies Agent 5's structured modification instructions to a PPTX file.

Supported actions:
  replace_text   — find original_text in a shape and replace with new_text
  append_bullets — add a list of bullet strings as new paragraphs to a shape
  append_text    — append a single text block as a new paragraph to a shape

All failures are soft — logged and returned, never crash the full run.
"""

import io
import logging
from copy import deepcopy
from pptx import Presentation
from lxml import etree

logger = logging.getLogger(__name__)


# ── Helpers ───────────────────────────────────────────────────────────────────

def _find_shape(slide, shape_name: str):
    """Find a shape on a slide by exact name."""
    for shape in slide.shapes:
        if shape.name == shape_name:
            return shape
    return None


def _fuzzy_contains(haystack: str, needle: str) -> bool:
    """Case-insensitive substring check, ignoring leading/trailing whitespace."""
    return needle.strip().lower() in haystack.strip().lower()


def _replace_paragraph_text(para, new_text: str) -> None:
    """
    Replace all text in a paragraph with new_text, preserving the formatting
    of the first run (font name, size, bold, italic, colour).

    Clears ALL non-property child elements (runs, field codes, line-breaks) so
    no old content survives alongside the new text.
    """
    NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
    p_elem = para._p

    # Capture first run's rPr (formatting) before we clear anything
    first_rpr = None
    for r in p_elem.findall(f"{{{NS}}}r"):
        rpr = r.find(f"{{{NS}}}rPr")
        if rpr is not None:
            first_rpr = deepcopy(rpr)
        break  # only need the first run's properties

    # Remove ALL child elements except pPr (paragraph-level properties:
    # bullet style, spacing, indent, alignment — these must be preserved)
    for child in list(p_elem):
        if child.tag != f"{{{NS}}}pPr":
            p_elem.remove(child)

    # Build one clean run with the new text
    r_elem = etree.SubElement(p_elem, f"{{{NS}}}r")
    if first_rpr is not None:
        r_elem.append(first_rpr)
    t_elem = etree.SubElement(r_elem, f"{{{NS}}}t")
    t_elem.text = new_text
    # Preserve leading/trailing spaces
    if new_text and (new_text[0] == " " or new_text[-1] == " "):
        t_elem.set("{http://www.w3.org/XML/1998/namespace}space", "preserve")


def _replace_text_in_shape(shape, original_text: str, new_text: str) -> bool:
    """
    Searches a shape's text frame for a paragraph containing original_text
    (fuzzy, case-insensitive) and replaces its full content with new_text.

    Returns True if a replacement was made, False otherwise.
    """
    if not shape.has_text_frame:
        return False

    tf = shape.text_frame

    # Try to find the paragraph that best matches original_text
    best_para = None
    best_ratio = 0.0
    needle_lower = original_text.strip().lower()

    for para in tf.paragraphs:
        para_text = para.text.strip().lower()
        if not para_text:
            continue
        # Exact or substring match
        if needle_lower in para_text or para_text in needle_lower:
            # Score by length similarity — prefer the closest match
            ratio = min(len(para_text), len(needle_lower)) / max(
                len(para_text), len(needle_lower), 1
            )
            if ratio > best_ratio:
                best_ratio = ratio
                best_para = para

    if best_para is None:
        return False

    _replace_paragraph_text(best_para, new_text)
    return True


def _append_paragraph_to_shape(shape, text: str, level: int = 1) -> None:
    """Add a new paragraph with the given text to the shape's text frame."""
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    new_para = tf.add_paragraph()
    new_para.text = text
    new_para.level = level


# ── Main apply function ───────────────────────────────────────────────────────

def apply_modifications(pptx_bytes: bytes, modifications: list) -> tuple:
    """
    Apply a list of modification dicts to the PPTX.

    Returns:
      (modified_pptx_bytes: bytes, applied: list, failed: list)
    """
    prs = Presentation(io.BytesIO(pptx_bytes))
    total_slides = len(prs.slides)

    applied = []
    failed = []

    # Sort: must_fix first, then should_fix, then nice_to_have
    _priority_order = {"must_fix": 0, "should_fix": 1, "nice_to_have": 2}
    sorted_mods = sorted(
        modifications,
        key=lambda m: _priority_order.get(m.get("priority", "nice_to_have"), 2),
    )

    for mod in sorted_mods:
        slide_index = mod.get("slide_index")
        shape_name = mod.get("shape_name", "")
        action = mod.get("action", "")
        original_text = mod.get("original_text", "")
        new_text = mod.get("new_text", "")
        bullets = mod.get("bullets", [])

        # ── Validate slide index ──────────────────────────────────────────────
        if slide_index is None or not isinstance(slide_index, int):
            failed.append({**mod, "error": "Missing or invalid slide_index"})
            continue
        if slide_index >= total_slides:
            failed.append({**mod, "error": f"slide_index {slide_index} out of range (total: {total_slides})"})
            continue

        slide = prs.slides[slide_index]

        # ── Find shape ────────────────────────────────────────────────────────
        shape = _find_shape(slide, shape_name)
        if shape is None:
            failed.append({**mod, "error": f"Shape '{shape_name}' not found on slide {slide_index}"})
            continue

        # ── Execute action ────────────────────────────────────────────────────
        try:
            if action == "replace_text":
                if not original_text or not new_text:
                    failed.append({**mod, "error": "replace_text requires both original_text and new_text"})
                    continue

                # Defence-in-depth: never replace text inside a large-font
                # non-placeholder shape.  These are design/callout elements whose
                # modified content would visually overlay real body text.
                if not shape.is_placeholder and shape.has_text_frame:
                    font_sizes = [
                        int(run.font.size.pt)
                        for para in shape.text_frame.paragraphs
                        for run in para.runs
                        if run.font.size
                    ]
                    if font_sizes and max(font_sizes) > 28:
                        failed.append({**mod, "error": (
                            f"Blocked: shape '{shape_name}' is a design element "
                            f"(max font {max(font_sizes)}pt). Edit manually."
                        )})
                        continue

                success = _replace_text_in_shape(shape, original_text, new_text)
                if success:
                    applied.append(mod)
                    logger.info(
                        "Applied replace_text on slide %d shape '%s'", slide_index, shape_name
                    )
                else:
                    failed.append({**mod, "error": f"original_text not found in shape '{shape_name}'"})

            elif action == "append_bullets":
                if not bullets:
                    failed.append({**mod, "error": "append_bullets requires a non-empty bullets list"})
                    continue
                for bullet in bullets:
                    _append_paragraph_to_shape(shape, bullet, level=1)
                applied.append(mod)
                logger.info(
                    "Applied append_bullets (%d items) on slide %d shape '%s'",
                    len(bullets), slide_index, shape_name,
                )

            elif action == "append_text":
                if not new_text:
                    failed.append({**mod, "error": "append_text requires new_text"})
                    continue
                _append_paragraph_to_shape(shape, new_text, level=0)
                applied.append(mod)
                logger.info(
                    "Applied append_text on slide %d shape '%s'", slide_index, shape_name
                )

            else:
                failed.append({**mod, "error": f"Unknown action '{action}'"})

        except Exception as exc:
            logger.error(
                "Error applying modification on slide %d shape '%s': %s",
                slide_index, shape_name, exc, exc_info=True,
            )
            failed.append({**mod, "error": str(exc)})

    # ── Serialise ─────────────────────────────────────────────────────────────
    output = io.BytesIO()
    prs.save(output)
    return output.getvalue(), applied, failed
