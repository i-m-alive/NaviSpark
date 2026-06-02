import io
import time
import logging
import threading
import boto3
import json
import base64
import re
from botocore.config import Config
from fastapi import HTTPException
from config import settings

logger = logging.getLogger(__name__)

# ── Per-thread token accumulator ─────────────────────────────────────────────
# Each agent's run() resets this at entry and reads it at exit, so parallel
# agents running in separate thread-pool threads never contaminate each other.
_token_accumulator = threading.local()


def reset_token_accumulator() -> None:
    """Reset the per-thread token counters to zero. Call at the start of each agent run."""
    _token_accumulator.input_tokens = 0
    _token_accumulator.output_tokens = 0


def get_accumulated_tokens() -> dict:
    """Return the tokens accumulated since the last reset on this thread."""
    inp = getattr(_token_accumulator, "input_tokens", 0)
    out = getattr(_token_accumulator, "output_tokens", 0)
    return {"input_tokens": inp, "output_tokens": out, "total_tokens": inp + out}


def _accumulate_tokens(usage: dict) -> None:
    _token_accumulator.input_tokens = getattr(_token_accumulator, "input_tokens", 0) + usage.get("input_tokens", 0)
    _token_accumulator.output_tokens = getattr(_token_accumulator, "output_tokens", 0) + usage.get("output_tokens", 0)

# boto3 client config — adaptive retry mode handles short throttle windows;
# our own _invoke_bedrock_with_retry handles sustained token-bucket exhaustion.
#
# Timeout rationale for Opus + 80-page PDFs:
#   connect_timeout : TCP handshake to AWS endpoint — 60 s is generous.
#   read_timeout    : Max time to wait between response bytes.
#                     Opus at ~100 tok/s × 16,000 output tokens = ~160 s just for
#                     generation, plus model warm-up and input processing on large
#                     documents can add another 60-120 s.
#                     600 s (10 min) gives a 3× safety margin for worst-case runs.
_BEDROCK_CONFIG = Config(
    retries={
        "max_attempts": 6,
        "mode": "adaptive",
    },
    read_timeout=600,
    connect_timeout=60,
)

# Delays (seconds) for our outer retry layer, on top of boto3's built-in retries.
# Used when the account's TPM quota is exhausted and we need to wait for refill.
_THROTTLE_RETRY_DELAYS = [20, 40, 60]


_PLACEHOLDERS = {"", "your-aws-access-key", "your-aws-secret-key"}


def get_bedrock_client():
    """
    Builds a Bedrock client using whichever credentials are available:
      1. Explicit keys in .env  — used when AWS_ACCESS_KEY_ID is set and not a placeholder.
         SSO temp creds require all three: ACCESS_KEY + SECRET_KEY + SESSION_TOKEN.
      2. AWS CLI profile / default chain — used when keys are absent/placeholder.
         Run `aws sso login` first, then set AWS_PROFILE in .env if non-default.
    """
    use_explicit = (
        settings.aws_access_key_id not in _PLACEHOLDERS
        and settings.aws_secret_access_key not in _PLACEHOLDERS
    )

    if use_explicit:
        kwargs = dict(
            service_name="bedrock-runtime",
            region_name=settings.aws_region,
            aws_access_key_id=settings.aws_access_key_id,
            aws_secret_access_key=settings.aws_secret_access_key,
            config=_BEDROCK_CONFIG,
        )
        if settings.aws_session_token:
            kwargs["aws_session_token"] = settings.aws_session_token
        return boto3.client(**kwargs)

    # Fall back to boto3 default credential chain (env vars, ~/.aws/credentials, SSO profile)
    session = boto3.Session(
        profile_name=settings.aws_profile or None,
        region_name=settings.aws_region,
    )
    return session.client("bedrock-runtime", config=_BEDROCK_CONFIG)


def _invoke_bedrock_with_retry(client, request_body: dict):
    """
    Wraps boto3 invoke_model with an outer retry layer for sustained ThrottlingExceptions.

    boto3's built-in adaptive retry handles short token-bucket exhaustion (a few seconds).
    This function handles the case where the TPM quota is depleted for longer — e.g.
    when multiple concurrent calls have already burned through the quota. It waits
    _THROTTLE_RETRY_DELAYS seconds between attempts (20s → 40s → 60s) giving Bedrock's
    token bucket time to refill before re-trying.
    """
    last_exc = None
    for attempt, delay in enumerate([0] + _THROTTLE_RETRY_DELAYS):
        if delay:
            logger.warning(
                "[BEDROCK] ThrottlingException — waiting %ds before retry %d/%d...",
                delay, attempt, len(_THROTTLE_RETRY_DELAYS),
            )
            time.sleep(delay)
            logger.info("[BEDROCK] Retrying Bedrock call now (attempt %d/%d)...",
                        attempt + 1, len(_THROTTLE_RETRY_DELAYS) + 1)
        try:
            return client.invoke_model(
                modelId=settings.bedrock_model_id,
                contentType="application/json",
                accept="application/json",
                body=json.dumps(request_body),
            )
        except Exception as exc:
            err = str(exc)
            if "ThrottlingException" in err or "Too many tokens" in err:
                last_exc = exc
                continue
            raise  # Non-throttling error — re-raise immediately
    logger.error("[BEDROCK] All %d retry attempts exhausted. Giving up.", len(_THROTTLE_RETRY_DELAYS) + 1)
    raise last_exc


def _clean_json_response(raw_text: str) -> str:
    """
    Extracts the first complete JSON object from the model response.

    Handles three common Opus / large-doc failure modes:
      1. Markdown code fences wrapping the JSON (```json ... ```)
      2. Preamble text before the opening brace ("Here is the JSON: {...}")
      3. Trailing text / commentary after the closing brace ({...} Note: ...)

    Strategy: strip fences first, then walk the string character-by-character
    tracking brace depth to find the exact span of the outermost JSON object.
    This is O(n) and handles arbitrarily nested structures.
    """
    # Step 1: strip markdown code fences
    text = raw_text.strip()
    text = re.sub(r'^```(?:json)?\s*', '', text)
    text = re.sub(r'\s*```$', '', text)
    text = text.strip()

    # Step 2: find the first '{' — everything before it is preamble
    start = text.find('{')
    if start == -1:
        return text  # no JSON object found; let caller handle the parse error

    # Step 3: walk forward tracking brace depth, skipping string contents
    depth = 0
    in_string = False
    escape_next = False

    for i, ch in enumerate(text[start:], start):
        if escape_next:
            escape_next = False
            continue
        if ch == '\\' and in_string:
            escape_next = True
            continue
        if ch == '"':
            in_string = not in_string
            continue
        if in_string:
            continue
        if ch == '{':
            depth += 1
        elif ch == '}':
            depth -= 1
            if depth == 0:
                return text[start : i + 1]   # exact JSON object, nothing extra

    # Brace depth never reached 0 — return from start to end (truncated JSON)
    return text[start:]


def _extract_pptx_text(pptx_bytes: bytes) -> str:
    """Extracts all text from a PPTX/PPT file, slide by slide."""
    from pptx import Presentation
    prs = Presentation(io.BytesIO(pptx_bytes))
    slides_text = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        if texts:
            slides_text.append(f"--- Slide {i} ---\n" + "\n".join(texts))
    return "\n\n".join(slides_text) if slides_text else "(No text content found in presentation)"


def invoke_agent_with_pdf(system_prompt: str, user_message: str, pdf_bytes: bytes, file_type: str = "pdf") -> dict:
    """
    Calls the configured Claude model on Bedrock with a document.
    - PDF: sent as a native base64 document block.
    - PPTX/PPT: text extracted slide-by-slide and prepended to the user message.
    """
    client = get_bedrock_client()

    if file_type in ("pptx", "ppt"):
        slide_content = _extract_pptx_text(pdf_bytes)
        content = [
            {
                "type": "text",
                "text": f"PROPOSAL CONTENT (PowerPoint presentation, {file_type.upper()}):\n\n{slide_content}\n\n{user_message}",
            }
        ]
    else:
        pdf_b64 = base64.standard_b64encode(pdf_bytes).decode("utf-8")
        content = [
            {
                "type": "document",
                "source": {
                    "type": "base64",
                    "media_type": "application/pdf",
                    "data": pdf_b64,
                },
            },
            {
                "type": "text",
                "text": user_message,
            },
        ]

    request_body = {
        "anthropic_version": "bedrock-2023-05-31",
        "max_tokens": 16000,
        "temperature": 0,
        "system": system_prompt,
        "messages": [
            {
                "role": "user",
                "content": content,
            }
        ],
    }

    try:
        response = _invoke_bedrock_with_retry(client, request_body)
    except Exception as e:
        raise HTTPException(status_code=502, detail=f"AWS Bedrock API call failed: {str(e)}")

    try:
        response_body = json.loads(response["body"].read())
        raw_text = response_body["content"][0]["text"]
        usage = response_body.get("usage", {})
        _accumulate_tokens(usage)
        logger.info(
            "[BEDROCK] tokens — input: %s  output: %s  (model context limit: 200,000)",
            usage.get("input_tokens", "?"),
            usage.get("output_tokens", "?"),
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to read Bedrock response: {str(e)}")

    cleaned = _clean_json_response(raw_text)
    try:
        return json.loads(cleaned)
    except json.JSONDecodeError as e:
        raise HTTPException(
            status_code=500,
            detail=f"Model returned invalid JSON. Error: {str(e)}. Raw (first 500): {raw_text[:500]}"
        )


def invoke_agent_text_only(system_prompt: str, user_message: str, max_tokens: int = 16000) -> dict:
    """Calls the configured Claude model on Bedrock with text input only (no PDF). Used by Agent 4 and Agent 5."""
    client = get_bedrock_client()

    request_body = {
        "anthropic_version": "bedrock-2023-05-31",
        "max_tokens": max_tokens,
        "temperature": 0,
        "system": system_prompt,
        "messages": [{"role": "user", "content": user_message}]
    }

    try:
        response = _invoke_bedrock_with_retry(client, request_body)
    except Exception as e:
        raise HTTPException(status_code=502, detail=f"AWS Bedrock API call failed: {str(e)}")

    try:
        response_body = json.loads(response["body"].read())
        raw_text = response_body["content"][0]["text"]
        usage = response_body.get("usage", {})
        _accumulate_tokens(usage)
        logger.info(
            "[BEDROCK] tokens — input: %s  output: %s  (model context limit: 200,000)",
            usage.get("input_tokens", "?"),
            usage.get("output_tokens", "?"),
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to read Bedrock response: {str(e)}")

    cleaned = _clean_json_response(raw_text)
    try:
        return json.loads(cleaned)
    except json.JSONDecodeError as e:
        raise HTTPException(
            status_code=500,
            detail=f"Model returned invalid JSON. Error: {str(e)}. Raw (first 500): {raw_text[:500]}"
        )


# ── Chunking: per-chunk summarisation ─────────────────────────────────────────

_CHUNK_SYSTEM_PROMPT = """You are a document analysis assistant specialising in business proposal analysis.
You extract structured information from sections of proposal documents.

Your output is consumed by specialist AI review agents. Accuracy is critical — every figure,
scope item, and commercial term you miss will be invisible to the agents that depend on this context.

RULES:
- Capture EVERY number, cost, estimate, timeline, and percentage — preserve exact values with units
- Capture EVERY scope item, deliverable, feature, and work item mentioned
- Capture ALL section headings and slide titles found in this section
- Note ALL risks, even passing mentions
- Note ALL commercial terms: pricing models, payment schedules, rate cards, cost categories
- Return ONLY valid JSON. No preamble. No explanation. No markdown fences. Start with { end with }."""


def summarize_chunk_pdf(
    chunk_bytes: bytes,
    chunk_index: int,
    page_start: int,
    page_end: int,
) -> dict:
    """
    Sends one PDF chunk (pages page_start–page_end) to Bedrock and returns
    a structured summary dict. Used by chunking_service for large PDFs.
    """
    client = get_bedrock_client()
    pdf_b64 = base64.standard_b64encode(chunk_bytes).decode("utf-8")

    user_message = (
        f"Analyse this section of a business proposal (pages {page_start}–{page_end}, "
        f"chunk index {chunk_index}).\n\n"
        f"Return ONLY this JSON structure — fill every field, use [] for empty arrays:\n"
        f'{{\n'
        f'  "chunk_index": {chunk_index},\n'
        f'  "page_range": {{"start": {page_start}, "end": {page_end}}},\n'
        f'  "summary": "<3-5 sentences describing what this section covers>",\n'
        f'  "sections_found": ["<heading or title>"],\n'
        f'  "key_claims": [{{"claim": "<assertion>", "location": "<section name>"}}],\n'
        f'  "figures_and_numbers": [{{"label": "<what this number represents>", "value": "<exact value with units>", "context": "<surrounding context>"}}],\n'
        f'  "scope_items": ["<deliverable, feature, or work item>"],\n'
        f'  "risks": [{{"risk": "<description>", "mitigation_present": true}}],\n'
        f'  "commercial_terms": ["<pricing model, payment term, rate, or cost category>"]\n'
        f'}}'
    )

    request_body = {
        "anthropic_version": "bedrock-2023-05-31",
        "max_tokens": 16000,
        "temperature": 0,
        "system": _CHUNK_SYSTEM_PROMPT,
        "messages": [
            {
                "role": "user",
                "content": [
                    {
                        "type": "document",
                        "source": {
                            "type": "base64",
                            "media_type": "application/pdf",
                            "data": pdf_b64,
                        },
                    },
                    {"type": "text", "text": user_message},
                ],
            }
        ],
    }

    try:
        response = _invoke_bedrock_with_retry(client, request_body)
    except Exception as e:
        raise HTTPException(
            status_code=502,
            detail=f"Bedrock chunk summarisation failed (chunk {chunk_index}, pages {page_start}-{page_end}): {str(e)}",
        )

    try:
        response_body = json.loads(response["body"].read())
        raw_text = response_body["content"][0]["text"]
    except Exception as e:
        raise HTTPException(
            status_code=500,
            detail=f"Failed to read Bedrock response for chunk {chunk_index}: {str(e)}",
        )

    cleaned = _clean_json_response(raw_text)
    try:
        return json.loads(cleaned)
    except json.JSONDecodeError:
        # Return a minimal fallback so the overall chunking pipeline keeps running
        return {
            "chunk_index": chunk_index,
            "page_range": {"start": page_start, "end": page_end},
            "summary": raw_text[:500] if raw_text else "Summarisation produced non-JSON output.",
            "sections_found": [],
            "key_claims": [],
            "figures_and_numbers": [],
            "scope_items": [],
            "risks": [],
            "commercial_terms": [],
        }


def summarize_chunk_pptx_text(
    slide_text: str,
    chunk_index: int,
    slide_start: int,
    slide_end: int,
) -> dict:
    """
    Sends one PPTX slide-text chunk (slides slide_start–slide_end) to Bedrock
    and returns a structured summary dict. Used by chunking_service for large PPTXs.
    """
    client = get_bedrock_client()

    user_message = (
        f"Analyse this section of a business proposal presentation "
        f"(slides {slide_start}–{slide_end}, chunk index {chunk_index}).\n\n"
        f"SLIDE CONTENT:\n{slide_text}\n\n"
        f"Return ONLY this JSON structure — fill every field, use [] for empty arrays:\n"
        f'{{\n'
        f'  "chunk_index": {chunk_index},\n'
        f'  "page_range": {{"start": {slide_start}, "end": {slide_end}}},\n'
        f'  "summary": "<3-5 sentences describing what these slides cover>",\n'
        f'  "sections_found": ["<slide title or section heading>"],\n'
        f'  "key_claims": [{{"claim": "<assertion>", "location": "<slide number or title>"}}],\n'
        f'  "figures_and_numbers": [{{"label": "<what this number represents>", "value": "<exact value with units>", "context": "<surrounding context>"}}],\n'
        f'  "scope_items": ["<deliverable, feature, or work item>"],\n'
        f'  "risks": [{{"risk": "<description>", "mitigation_present": true}}],\n'
        f'  "commercial_terms": ["<pricing model, payment term, rate, or cost category>"]\n'
        f'}}'
    )

    request_body = {
        "anthropic_version": "bedrock-2023-05-31",
        "max_tokens": 16000,
        "temperature": 0,
        "system": _CHUNK_SYSTEM_PROMPT,
        "messages": [{"role": "user", "content": user_message}],
    }

    try:
        response = _invoke_bedrock_with_retry(client, request_body)
    except Exception as e:
        raise HTTPException(
            status_code=502,
            detail=f"Bedrock chunk summarisation failed (chunk {chunk_index}, slides {slide_start}-{slide_end}): {str(e)}",
        )

    try:
        response_body = json.loads(response["body"].read())
        raw_text = response_body["content"][0]["text"]
    except Exception as e:
        raise HTTPException(
            status_code=500,
            detail=f"Failed to read Bedrock response for chunk {chunk_index}: {str(e)}",
        )

    cleaned = _clean_json_response(raw_text)
    try:
        return json.loads(cleaned)
    except json.JSONDecodeError:
        return {
            "chunk_index": chunk_index,
            "page_range": {"start": slide_start, "end": slide_end},
            "summary": raw_text[:500] if raw_text else "Summarisation produced non-JSON output.",
            "sections_found": [],
            "key_claims": [],
            "figures_and_numbers": [],
            "scope_items": [],
            "risks": [],
            "commercial_terms": [],
        }


def invoke_agent_with_context_json(
    system_prompt: str,
    user_message: str,
    context_json: str,
    max_tokens: int = 16000,
) -> dict:
    """
    Calls an agent on Bedrock using pre-processed document context instead of a raw file.
    Used for large documents (>CHUNK_THRESHOLD pages) where the file has been split,
    summarised per-chunk, and merged into a unified context JSON by chunking_service.

    The context JSON is prepended to the user message so the agent reads the full
    document context before its task instructions.
    """
    combined_message = (
        "PROPOSAL CONTEXT (pre-processed from a large document via chunking pipeline — "
        "treat this structured summary as the full proposal):\n\n"
        f"{context_json}\n\n"
        "---\n\n"
        f"{user_message}"
    )
    return invoke_agent_text_only(system_prompt, combined_message, max_tokens)
